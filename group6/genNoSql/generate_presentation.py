from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Helper Functions ---
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_content_slide(title, content_points, image_placeholder_text=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Add Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Clear default text
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 0
            p.space_after = Pt(10)

        # Add Image Placeholder if requested
        if image_placeholder_text:
            left = Inches(5.5)
            top = Inches(2.0)
            width = Inches(4.0)
            height = Inches(3.0)
            
            # Draw a box for the placeholder
            shape = slide.shapes.add_shape(
                1, left, top, width, height  # 1 is MSO_SHAPE.RECTANGLE
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
            shape.text = f"[PLACEHOLDER]\n{image_placeholder_text}\n(Refer to sql_gen_paper.pdf)"

    # --- Slide 1: Title ---
    add_title_slide(
        "Text-to-NoSQL: NL to MongoDB Translation",
        "Capstone Stage 2 Final Report\nSMART Framework + RAG Optimization\n\nJanuary 1, 2026"
    )

    # --- Slide 2: Project Aim ---
    add_content_slide(
        "Project Aim: NL to NoSQL",
        [
            "Phase: Capstone Stage 2",
            "Goal: Translate Natural Language Questions into executable MongoDB queries.",
            "Target Baseline: 83% Execution Success Rate.",
            "Dataset: Spider Dev (1,034 MongoDB samples).",
            "Models Tested: Llama 3.2 (3B), Llama 3.1 (8B), Mistral 7B.",
            "Challenge: Diverse database schemas (~14 distinct types) and complex aggregations."
        ],
        "Insert Project Overview Image (if available)"
    )

    # --- Slide 3: SMART Architecture ---
    add_content_slide(
        "SMART Architecture Implementation",
        [
            "S: Semantically Enhanced Retrieval (RAG)",
            "M: Multi-turn Refinement",
            "A: Abstract Syntax Tree Validation (Execution Optimizer)",
            "R: Ranking & Selection",
            "T: Translation to Final Query",
            "",
            "Key Innovation: Integrated RAG to 'teach' the model syntax patterns dynamically."
        ],
        "Insert Architecture Diagram from sql_gen_paper.pdf\n(Start of paper, showing pipeline)"
    )

    # --- Slide 4: Methodology - The "Balanced" Strategy ---
    add_content_slide(
        "Methodology: Prompt Engineering & RAG",
        [
            "Evolution of Prompt Strategy:",
            "❌ Verbose (~60 lines): Overwhelmed the model (50% success).",
            "❌ Minimal (~10 lines): Led to hallucinations (30% success).",
            "✅ Balanced (~25 lines): Strict syntax rules + RAG examples.",
            "",
            "RAG Implementation:",
            "- Indexed 100 training examples (indices 0-99).",
            "- Retriever: 'all-MiniLM-L6-v2' embeddings.",
            "- Refiner: Appends 3 most similar verified queries."
        ],
        "Insert Prompt/RAG Flow Diagram from sql_gen_paper.pdf"
    )

    # --- Slide 5: The "Zero to Hero" Journey ---
    add_content_slide(
        "Phase 1: Critical Bug Fixes (0% -> 60%)",
        [
            "Initial State: 0% Success (System Crashes).",
            "",
            "1. Foundational Fixes:",
            "   - Fixed MongoDB 'Phantom Connection' bug.",
            "   - Fixed Schema Prediction format mismatch.",
            "",
            "2. The 'JavaScript Crisis':",
            "   - Issue: Llama generates JS 'true/null'; Python crashes.",
            "   - Fix: Context-aware Regex Normalization.",
            "   - Result: Boolean errors reduced from 23 to 0."
        ]
    )

    # --- Slide 6: Model Comparison ---
    add_content_slide(
        "Model Comparison: Llama vs Mistral",
        [
            "Tested on 50 Diverse Samples (14 databases):",
            "",
            "1. Llama 3.2 (3B):",
            "   - Stability: 52%",
            "   - Strength: Schema guessing (48% match).",
            "",
            "2. Mistral 7B:",
            "   - Stability: 68%",
            "   - Weakness: Syntax errors (unbalanced brackets).",
            "",
            "3. Llama 3.1 (8B) [WINNER]:",
            "   - Stability: 98% (Production Grade)",
            "   - Weakness: Verbose schema output (36% match)."
        ]
    )

    # --- Slide 7: Large Scale Results ---
    add_content_slide(
        "Validation: Large Scale Testing",
        [
            "Sample Size: 139 Queries (Full Dataset subset).",
            "Model: Llama 3.1 8B.",
            "",
            "Results:",
            "✅ Pipeline Success: 79/139 (57%)",
            "✅ Stability: Consistent with small-scale tests.",
            "",
            "Conclusion:",
            "System is stable and proof-of-concept is validated.",
            "57% is a strong baseline for diverse, unseen schemas."
        ]
    )

    # --- Slide 8: Key Metrics Summary ---
    add_content_slide(
        "Key Metrics & Improvements",
        [
            "Pipeline Stability:",
            "   0% (Start) -> 98% (Final/8B)",
            "",
            "Homogeneous Accuracy (1 DB):",
            "   70% Success.",
            "",
            "Diverse Accuracy (14 DBs):",
            "   57% Success (Stable).",
            "",
            "Boolean/Null Errors:",
            "   100% Eliminated."
        ],
        "Insert Performance Chart from sql_gen_paper.pdf"
    )

    # --- Slide 9: Future Roadmap (to 83%) ---
    add_content_slide(
        "Roadmap to 83% Baseline",
        [
            "Phase 1: Stratified Training (Week 1)",
            "- Expand RAG from 100 -> 300 examples.",
            "- Ensure all 14 database types are represented.",
            "",
            "Phase 2: Operational Specialization (Week 2)",
            "- Split RAG into 'Aggregate' vs 'Find' stores.",
            "- Teach model specifically for complex GroupBy.",
            "",
            "Phase 3: Edge Case Hardening",
            "- Add 'Self-Correction' loop for invalid syntax."
        ]
    )

    # --- Slide 10: Conclusion ---
    add_content_slide(
        "Conclusion",
        [
            "1. Successfully established a functional NL-to-NoSQL pipeline.",
            "2. Proved the 'Balanced Prompt + RAG' strategy works.",
            "3. Identified Llama 3.1 8B as the ideal local model.",
            "4. Clear, data-driven path to reach the 83% target.",
            "",
            "Ready for Phase 3: Advanced Optimization."
        ]
    )

    # Save Presentation
    output_file = "Capstone_Stage2_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
